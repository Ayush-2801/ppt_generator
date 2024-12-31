from langchain.llms import Ollama
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from pptx import Presentation
from pptx.util import Pt

llm = Ollama(model="gemma:2b")
prompt_template = PromptTemplate(
    input_variables=["topic"],
    template="Create a detailed PowerPoint slide content on the topic: {topic}. Include key points and a brief explanation. Break it down into sections so that it can be distributed across multiple slides."
)

chain = LLMChain(llm=llm, prompt=prompt_template)

def generate_content(topic):
    content = chain.run({"topic": topic})
    return content


def divide_content(content, num_slides):
    sections = content.split("\n")
    sections_per_slide = len(sections) // num_slides
    remainder = len(sections) % num_slides
    slides_content = []
    start_idx = 0
    for i in range(num_slides):
        end_idx = start_idx + sections_per_slide + (1 if i < remainder else 0)
        slide_section = "\n".join(sections[start_idx:end_idx])
        slides_content.append(slide_section)
        start_idx = end_idx
    return slides_content

def create_presentation(topic, content, num_slides, name):
    slides_content = divide_content(content, num_slides)
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Presentation on {topic}"
    subtitle.text = f"Created by:- {name}"
    for i, slide_content in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{topic}"
        content_box = slide.shapes.placeholders[1]
        content_box.text = slide_content
        for paragraph in content_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
    prs.save(f"{topic}_presentation_{num_slides}_slides.pptx")
    print(f"Presentation on '{topic}' with {num_slides} slides has been created successfully!")

def add_bullet_points(slide, content):
    textbox = slide.shapes.placeholders[1]
    text_frame = textbox.text_frame
    text_frame.clear()
    for point in content.split("\n"):
        p = text_frame.add_paragraph()
        p.text = point
        for run in p.runs:
            run.font.size = Pt(20)

def create_presentation_with_bullets(topic, content, num_slides, name):
    slides_content = divide_content(content, num_slides)
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Presentation on {topic}"
    subtitle.text = f"Created by:- {name}"
    for i, slide_content in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{topic}"
        add_bullet_points(slide, slide_content)
    prs.save(f"{topic}_presentation_with_bullets_{num_slides}_slides.pptx")
    print(f"Presentation with bullet points on '{topic}' and {num_slides} slides has been created successfully!")

name = input("Enter your name: ").capitalize()
topic = input("Enter the topic for your presentation: ")
num_slides = int(input("How many slides do you want in the presentation? "))
topic = topic.capitalize()
content = generate_content(topic)
choice = input("Do you want to create a presentation with bullet points? (yes/no): ").lower()
if choice == 'yes':
    create_presentation_with_bullets(topic, content, num_slides, name)
else:
    create_presentation(topic, content, num_slides, name)