from flask import Flask, render_template, jsonify, request, send_from_directory
from langchain.llms import Ollama
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from pptx import Presentation
from pptx.util import Pt
import os
import re

app = Flask(__name__)

llm = Ollama(model="gemma:2b")
prompt_template = PromptTemplate(
    input_variables=["topic"],
    template="Create a detailed PowerPoint slide content on the topic: {topic}. Include key points and a brief explanation. Break it down into sections so that it can be distributed across multiple slides."
)
chain = LLMChain(llm=llm, prompt=prompt_template)
    
def generate_content(topic):
    try:
        content = chain.run({"topic": topic})
        return content
    except Exception as e:
        print(f"Error generating content: {e}")
        return ""

def divide_content(content, num_slides):
    sections = content.split("\n")
    sections_per_slide = len(sections) // num_slides
    remainder = len(sections) % num_slides
    slides_content = []
    start_idx = 0
    for i in range(num_slides):
        end_idx = start_idx + sections_per_slide + (1 if i < remainder else 0)
        slide_section = "\n".join(sections[start_idx:end_idx])
        slides_content.append(slide_section)
        start_idx = end_idx
    return slides_content

def add_bullet_points(slide, content):
    textbox = slide.shapes.placeholders[1]
    text_frame = textbox.text_frame
    text_frame.clear()
    for point in content.split("\n"):
        p = text_frame.add_paragraph()
        p.text = point
        for run in p.runs:
            run.font.size = Pt(20)

def create_presentation_with_bullets(topic, content, num_slides, name, clean):
    slides_content = divide_content(content, num_slides)
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Presentation on {topic}"
    subtitle.text = f"Created by:- {name}"
    for i, slide_content in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{topic}"
        add_bullet_points(slide, slide_content)
    prs.save(f"Powerpoint/{clean}.pptx")
    print(f"Presentation with bullet points on '{topic}' and {num_slides} slides has been created successfully!")
    return True

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/create-ppt', methods=['POST'])
def create_ppt():
    if request.form['name']!="" and request.form['title']!="" and request.form['numofslide']!="":
        directory = os.listdir('Powerpoint')
        if len(directory) > 0:
            for file in directory:
                file_path = os.path.join('Powerpoint', file)
                if os.path.isfile(file_path):
                    os.remove(file_path)
        name = request.form['name'].capitalize()
        topic = request.form['title'].capitalize()
        clean = re.sub(r'[^a-zA-Z0-9]', '', topic)
        number_of_slides = int(request.form['numofslide'])
        content = generate_content(clean)
        created = create_presentation_with_bullets(topic=topic, content=content, num_slides=number_of_slides, name=name, clean=clean)
        if created:            
            return jsonify({'message': "PPT is created", "topic": clean})
        
    else:
        messageError = "Please provide all the fields."
        return jsonify({"messageError":messageError})
    
@app.route('/Powerpoint/<filename>')
def download_ppt(filename):
    return send_from_directory('Powerpoint', filename)

if __name__=='__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)